# 校验生成的 pptx: zip 完整性 / 全部 XML well-formed / 关系与部件一致 / python-pptx 可打开
# 用法: python validate.py gen/out_all.pptx gen/out_dark.pptx gen/out_outline.pptx gen/out_edit.pptx
import re
import sys
import zipfile
import json
from xml.etree import ElementTree as ET

from pptx import Presentation

R_ID = re.compile(rb'r:(?:embed|id|link)="(rId\d+)"')

def check_rels(z, names):
    """每页 slideN.xml 引用的 rId 必须存在于其 rels, 且 rels 的 Target 必须存在。"""
    import posixpath
    errors = []
    rels_names = [n for n in names if re.match(r'ppt/slides/_rels/slide\d+\.xml\.rels$', n)]
    for rels_name in rels_names:
        slide_name = rels_name.replace('_rels/', '').replace('.rels', '')
        rels_root = ET.fromstring(z.read(rels_name))
        rid_map = {}
        for rel in rels_root:
            rid = rel.get('Id')
            target = rel.get('Target')
            rid_map[rid] = target
            # 解析相对目标
            full = posixpath.normpath(posixpath.join(posixpath.dirname(slide_name), target))
            if full not in names:
                errors.append(f'{rels_name}: 目标不存在 {target}')
        used = set(R_ID.findall(z.read(slide_name)))
        for rid in used:
            if rid.decode() not in rid_map:
                errors.append(f'{slide_name}: 引用了未定义的 {rid.decode()}')
    return errors

def check_content_types(z, names):
    root = ET.fromstring(z.read('[Content_Types].xml'))
    ns = {'ct': 'http://schemas.openxmlformats.org/package/2006/content-types'}
    defaults = {e.get('Extension').lower() for e in root.findall('ct:Default', ns)}
    overrides = {e.get('PartName') for e in root.findall('ct:Override', ns)}
    errors = []
    for n in names:
        ext = n.rsplit('.', 1)[-1].lower()
        part = '/' + n
        if ext not in defaults and part not in overrides and not n.endswith('/'):
            errors.append(f'未登记的部件: {n}')
    return errors

def validate(path):
    print(f'== {path} ==')
    z = zipfile.ZipFile(path)
    assert z.testzip() is None, 'zip CRC 损坏'
    names = set(z.namelist())

    for n in names:
        if n.endswith('.xml') or n.endswith('.rels'):
            ET.fromstring(z.read(n))  # well-formed 校验
    print(f'  XML well-formed: OK ({len([n for n in names if n.endswith((".xml", ".rels"))])} 个部件)')

    errors = check_rels(z, names) + check_content_types(z, names)
    for e in errors:
        print('  RELS/CT FAIL:', e)
    assert not errors

    deck = json.loads(z.read('docProps/deck.json'))
    assert 'slides' in deck and len(deck['slides']) > 0

    prs = Presentation(path)
    n_slides = len(prs.slides)
    assert n_slides == len(deck['slides']), f'页数不一致: python-pptx {n_slides} vs deck {len(deck["slides"])}'
    print(f'  python-pptx: 打开成功, {n_slides} 页; 内嵌 deck.json 页数一致')

    # 统计关键部件
    media = [n for n in names if n.startswith('ppt/media/')]
    charts = [n for n in names if re.match(r'ppt/charts/chart\d+\.xml$', n)]
    notes = [n for n in names if re.match(r'ppt/notesSlides/notesSlide\d+\.xml$', n)]
    n_chart_slides = sum(1 for s in deck['slides'] if s.get('layout') == 'chart')
    n_notes = sum(1 for s in deck['slides'] if s.get('notes', '').strip())
    assert len(charts) == n_chart_slides, f'图表部件 {len(charts)} != 版式图表页 {n_chart_slides}'
    assert len(notes) == n_notes, f'备注部件 {len(notes)} != 备注页 {n_notes}'
    print(f'  media={len(media)} charts={len(charts)} notes={len(notes)}')

    # 抽取部分文本核对
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    joined = ' '.join(texts)
    return deck, joined, n_slides

def main():
    results = {}
    for path in sys.argv[1:]:
        deck, joined, n = validate(path)
        results[path] = (deck, joined)
    print('ALL VALIDATE OK')
    return results

if __name__ == '__main__':
    main()
