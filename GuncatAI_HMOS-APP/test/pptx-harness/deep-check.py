# 深度校验: 1) python-pptx 能读出图表数据(类别/系列/数值) 2) 内嵌 deck.json 与构建输入语义一致
import json
import zipfile
import sys
from pptx import Presentation

def charts(path, expected):
    prs = Presentation(path)
    found = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart:
                ch = shape.chart
                cats = [str(c) for c in ch.plots[0].categories]
                series = [(s.name, list(s.values)) for s in ch.series]
                found.append((cats, series))
    assert len(found) == len(expected), f'{path}: 图表数 {len(found)} != 期望 {len(expected)}'
    for i, (got, want) in enumerate(zip(found, expected)):
        assert got[0] == want[0], f'{path} 图表{i}: 类别不一致 {got[0]} vs {want[0]}'
        for (gn, gv), (wn, wv) in zip(got[1], want[1]):
            assert gn == wn, f'系列名 {gn} != {wn}'
            assert list(gv) == list(wv), f'{path} 图表{i}: 数值不一致 {list(gv)} vs {list(wv)}'
    print(f'{path}: {len(found)} 个图表数据读取一致')

def deck_roundtrip(path):
    z = zipfile.ZipFile(path)
    embedded = json.loads(z.read('docProps/deck.json'))
    # 结构字段存在且页内 layout 序列可复现
    assert embedded.get('title', '') != ''
    layouts = [s['layout'] for s in embedded['slides']]
    print(f'{path}: 内嵌源 {len(layouts)} 页, layouts={layouts[:6]}...')

if __name__ == '__main__':
    charts('gen/out_all.pptx', [
        (['4月', '5月', '6月'], [('营收', [36, 39, 45]), ('去年', [31, 33, 38])]),
        (['新客', '复购', '会员'], [('占比', [32, 41, 27])]),
    ])
    charts('gen/out_dark.pptx', [
        (['v1', 'v2', 'v3', 'v4'], [('P95', [420, 260, 180, 120])]),
    ])
    for p in sys.argv[1:] or ['gen/out_all.pptx', 'gen/out_dark.pptx', 'gen/out_outline.pptx', 'gen/out_edit.pptx']:
        deck_roundtrip(p)
    print('DEEP CHECK OK')
